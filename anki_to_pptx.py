#!/usr/bin/env python3
from __future__ import annotations

import argparse
import csv
import html
import io
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Dict, List

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

DEFAULT_MEDIA_DIR = Path("/Users/m/Library/Application Support/Anki2/test/collection.media")

# =========================
# CONSTANTS
# =========================

# Left/right deck slide size
LEFT_RIGHT_SLIDE_WIDTH_IN = 13.333
LEFT_RIGHT_SLIDE_HEIGHT_IN = 5.0

# Top/down deck slide size
TOP_DOWN_SLIDE_WIDTH_IN = 8.0
TOP_DOWN_SLIDE_HEIGHT_IN = 7.5

# 4-up printable page deck slide size (US Letter portrait)
PRINT4_PAGE_WIDTH_IN = 8.5
PRINT4_PAGE_HEIGHT_IN = 11.0

# Shared spacing
MARGIN_IN = 0.35
GUTTER_IN = 0.20
TEXT_AFTER_EN_PT = 10

# Left/right layout fonts
LEFT_RIGHT_EN_FONT_PT = 24
LEFT_RIGHT_ZH_FONT_PT = LEFT_RIGHT_EN_FONT_PT + 4

# Top/down layout fonts
TOP_DOWN_EN_FONT_PT = 20
TOP_DOWN_ZH_FONT_PT = TOP_DOWN_EN_FONT_PT + 3

# Top/down image size
TOP_DOWN_IMAGE_WIDTH_IN = 8.0
TOP_DOWN_IMAGE_HEIGHT_IN = 6.0
TOP_DOWN_TEXT_GAP_IN = 0.00
TOP_DOWN_TEXT_LEFT_INSET_FRACTION = 0.02
TOP_DOWN_TEXT_RIGHT_INSET_FRACTION = 0.02

# Image trimming
CROP_LEFT_PX = 14
CROP_RIGHT_PX = 14
CROP_TOP_PX = 0
CROP_BOTTOM_PX = 0

# 4-up page layout
PRINT4_PAGE_MARGIN_TOP_IN = 0.30
PRINT4_PAGE_MARGIN_BOTTOM_IN = 0.35
PRINT4_PAGE_MARGIN_LEFT_IN = 0.22
PRINT4_PAGE_MARGIN_RIGHT_IN = 0.22
PRINT4_CARD_GAP_X_IN = 0.15
PRINT4_CARD_GAP_Y_IN = 0.22
PRINT4_CARD_MARGIN_IN = 0.09
PRINT4_TEXT_GAP_IN = 0.05

# 4-up text formatting
PRINT4_EN_FONT_PT = 20
PRINT4_ZH_FONT_PT = 23
PRINT4_TEXT_LEFT_INSET_FRACTION = 0.02
PRINT4_TEXT_RIGHT_INSET_FRACTION = 0.02
PRINT4_EN_AFTER_PT = 4

# 4-up image sizing within each card
PRINT4_IMAGE_HEIGHT_FRACTION = 0.60

# 4-up fit strictness
PRINT4_BOTTOM_RESERVE_FRACTION = 0.10
PRINT4_FIT_WIDTH_SAFETY_FRACTION = 0.92
PRINT4_FIT_HEIGHT_SAFETY_FRACTION = 0.84
PRINT4_MEASURE_LINE_GAP_FRACTION = 0.36

# PDF conversion
AUTO_CONVERT_TO_PDF = True

# Auto-fit safeguards
MIN_FONT_PT = 8
FONT_SHRINK_STEP_PT = 1

# Overlap cleanup
MAX_OVERLAP_SEGMENTS_TO_CHECK = 3


def clean_text(text: str) -> str:
    text = html.unescape(text or "")
    text = text.replace("<br>", " ").replace("<br/>", " ").replace("<br />", " ")
    text = re.sub(r"<[^>]+>", "", text)
    text = text.replace("\r\n", " ").replace("\r", " ").replace("\n", " ")
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def parse_anki_export(input_path: Path) -> List[Dict[str, str]]:
    raw = input_path.read_text(encoding="utf-8-sig")
    lines = raw.splitlines()

    while lines and lines[0].startswith("#"):
        lines.pop(0)

    body = "\n".join(lines)
    reader = csv.reader(io.StringIO(body), delimiter="\t", quotechar='"')
    parsed_rows: List[Dict[str, str]] = []

    for row in reader:
        if not row:
            continue

        cell = next((c for c in row if c.strip()), "")
        if not cell:
            continue

        img_match = re.search(r'<img\s+src="([^"]+)"', cell, re.IGNORECASE)
        img_name = img_match.group(1).strip() if img_match else ""

        div_matches = list(re.finditer(r"<div[^>]*>(.*?)</div>", cell, re.DOTALL | re.IGNORECASE))
        zh_raw = ""
        audio_name = ""
        en_section = cell

        if div_matches:
            first_end = div_matches[0].end()
            en_section = cell[first_end:]

            if len(div_matches) >= 2:
                zh_raw = div_matches[1].group(1)
                en_section = cell[first_end:div_matches[1].start()]

            if len(div_matches) >= 3:
                audio_content = div_matches[2].group(1)
                audio_match = re.search(r"\[sound:([^\]]+)\]", audio_content)
                if audio_match:
                    audio_name = audio_match.group(1).strip()

        parsed_rows.append(
            {
                "En": clean_text(en_section),
                "Zh": clean_text(zh_raw),
                "ImgBaseName": img_name,
                "AudioBaseName": audio_name,
            }
        )

    return parsed_rows


def export_sanitized_csv(rows: List[Dict[str, str]], output_path: Path) -> None:
    with output_path.open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=["En", "Zh", "ImgBaseName", "AudioBaseName"])
        writer.writeheader()
        writer.writerows(rows)


def split_into_segments(text: str) -> List[str]:
    text = re.sub(r"\s+", " ", (text or "")).strip()
    if not text:
        return []

    # Split on sentence-ending punctuation, preserving punctuation.
    parts = re.split(r"(?<=[.!?。！？])\s+", text)
    parts = [p.strip() for p in parts if p.strip()]

    # Fallback: if no sentence punctuation found, split on spaced em dashes / semicolons.
    if len(parts) == 1:
        parts = re.split(r"\s+(?=--|—)|(?<=;)\s+|(?<=；)\s+", text)
        parts = [p.strip() for p in parts if p.strip()]

    return parts


def normalize_segment(text: str) -> str:
    text = text.lower().strip()
    text = re.sub(r"[“”\"'‘’]", "", text)
    text = re.sub(r"\s+", " ", text)
    return text


def remove_trailing_overlap(curr_text: str, next_text: str) -> str:
    curr_segments = split_into_segments(curr_text)
    next_segments = split_into_segments(next_text)

    if not curr_segments or not next_segments:
        return curr_text

    max_k = min(len(curr_segments), len(next_segments), MAX_OVERLAP_SEGMENTS_TO_CHECK)

    for k in range(max_k, 0, -1):
        curr_suffix = [normalize_segment(x) for x in curr_segments[-k:]]
        next_prefix = [normalize_segment(x) for x in next_segments[:k]]
        if curr_suffix == next_prefix:
            trimmed = curr_segments[:-k]
            return " ".join(trimmed).strip()

    return curr_text


def dedupe_neighbor_overlaps(rows: List[Dict[str, str]]) -> List[Dict[str, str]]:
    if not rows:
        return rows

    cleaned = [dict(r) for r in rows]

    for i in range(len(cleaned) - 1):
        cleaned[i]["En"] = remove_trailing_overlap(cleaned[i]["En"], cleaned[i + 1]["En"])
        cleaned[i]["Zh"] = remove_trailing_overlap(cleaned[i]["Zh"], cleaned[i + 1]["Zh"])

    return cleaned


def make_trimmed_image(image_path: Path) -> Path:
    with Image.open(image_path) as img:
        width, height = img.size

        left = max(0, min(CROP_LEFT_PX, width - 1))
        top = max(0, min(CROP_TOP_PX, height - 1))
        right = max(left + 1, min(width - CROP_RIGHT_PX, width))
        bottom = max(top + 1, min(height - CROP_BOTTOM_PX, height))

        cropped = img.crop((left, top, right, bottom))

        suffix = image_path.suffix if image_path.suffix else ".png"
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        tmp_path = Path(tmp.name)
        tmp.close()

        save_format = img.format if img.format else None
        cropped.save(tmp_path, format=save_format)
        return tmp_path


def get_prepared_image_path(image_path: Path | None) -> Path | None:
    if not image_path or not image_path.exists():
        return None
    return make_trimmed_image(image_path)


def points_to_pixels(pt: int, dpi: int = 96) -> int:
    return max(1, int(pt * dpi / 72))


def wrap_text_to_width(text: str, font, max_width_px: int, draw) -> List[str]:
    words = text.split()
    if not words:
        return [""]

    lines: List[str] = []
    current = words[0]

    for word in words[1:]:
        trial = current + " " + word
        bbox = draw.textbbox((0, 0), trial, font=font)
        trial_width = bbox[2] - bbox[0]

        if trial_width <= max_width_px:
            current = trial
        else:
            lines.append(current)
            current = word

    lines.append(current)
    return lines


def _load_measure_font(font_size_pt: int):
    px = points_to_pixels(font_size_pt)
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial Unicode.ttf",
        "/Library/Fonts/Arial.ttf",
    ]

    for candidate in candidates:
        try:
            return ImageFont.truetype(candidate, px)
        except Exception:
            pass

    return ImageFont.load_default()


def measure_wrapped_text_height(text: str, font_size_pt: int, max_width_px: int) -> int:
    img = Image.new("RGB", (10, 10), "white")
    draw = ImageDraw.Draw(img)
    font = _load_measure_font(font_size_pt)

    lines = wrap_text_to_width(text, font, max_width_px, draw)

    bbox = draw.textbbox((0, 0), "Ag", font=font)
    line_height = bbox[3] - bbox[1]
    line_gap = max(1, int(line_height * PRINT4_MEASURE_LINE_GAP_FRACTION))

    total_height = len(lines) * line_height
    if len(lines) > 1:
        total_height += (len(lines) - 1) * line_gap

    return total_height


def fit_font_sizes_for_two_paragraphs(
    en_text: str,
    zh_text: str,
    box_width_emu: int,
    box_height_emu: int,
    start_en_pt: int,
    start_zh_pt: int,
    paragraph_gap_pt: int,
) -> tuple[int, int]:
    emu_per_inch = 914400
    dpi = 96

    safe_width_emu = int(box_width_emu * PRINT4_FIT_WIDTH_SAFETY_FRACTION)
    safe_height_emu = int(box_height_emu * PRINT4_FIT_HEIGHT_SAFETY_FRACTION)

    box_width_px = max(1, int(safe_width_emu / emu_per_inch * dpi))
    box_height_px = max(1, int(safe_height_emu / emu_per_inch * dpi))
    paragraph_gap_px = points_to_pixels(paragraph_gap_pt, dpi=dpi)

    en_pt = start_en_pt
    zh_pt = start_zh_pt

    while en_pt >= MIN_FONT_PT and zh_pt >= MIN_FONT_PT:
        en_h = measure_wrapped_text_height(en_text, en_pt, box_width_px)
        zh_h = measure_wrapped_text_height(zh_text, zh_pt, box_width_px)

        total_h = en_h + paragraph_gap_px + zh_h
        if total_h <= box_height_px:
            return en_pt, zh_pt

        en_pt -= FONT_SHRINK_STEP_PT
        zh_pt -= FONT_SHRINK_STEP_PT

    return max(en_pt, MIN_FONT_PT), max(zh_pt, MIN_FONT_PT)


def add_slide_left_right(prs: Presentation, en_text: str, zh_text: str, image_path: Path | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    margin = Inches(MARGIN_IN)
    gutter = Inches(GUTTER_IN)

    left_x = margin
    left_y = margin
    left_w = slide_width / 2 - margin - gutter / 2
    left_h = slide_height - 2 * margin

    right_x = slide_width / 2 + gutter / 2
    right_y = margin
    right_w = slide_width / 2 - margin - gutter / 2
    right_h = slide_height - 2 * margin

    prepared_image = get_prepared_image_path(image_path)

    if prepared_image:
        pic = slide.shapes.add_picture(str(prepared_image), 0, 0)
        iw = pic.width
        ih = pic.height
        scale = min(left_w / iw, left_h / ih)
        new_w = int(iw * scale)
        new_h = int(ih * scale)
        pic.width = new_w
        pic.height = new_h
        pic.left = int(left_x + (left_w - new_w) / 2)
        pic.top = int(left_y + (left_h - new_h) / 2)
        content_top = pic.top
    else:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            int(left_x),
            int(left_y),
            int(left_w),
            int(left_h),
        )
        box.text_frame.text = "Image not found"
        box.fill.background()
        box.line.fill.background()
        content_top = left_y

    text_box = slide.shapes.add_textbox(
        int(right_x),
        int(content_top),
        int(right_w),
        int(right_h - (content_top - right_y)),
    )
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p1 = tf.paragraphs[0]
    p1.text = en_text
    p1.font.size = Pt(LEFT_RIGHT_EN_FONT_PT)
    p1.space_after = Pt(TEXT_AFTER_EN_PT)

    p2 = tf.add_paragraph()
    p2.text = zh_text
    p2.font.size = Pt(LEFT_RIGHT_ZH_FONT_PT)


def add_slide_top_image_text_below(prs: Presentation, en_text: str, zh_text: str, image_path: Path | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    bottom_margin = Inches(MARGIN_IN)
    text_gap = Inches(TOP_DOWN_TEXT_GAP_IN)

    img_w = Inches(TOP_DOWN_IMAGE_WIDTH_IN)
    img_h = Inches(TOP_DOWN_IMAGE_HEIGHT_IN)

    img_left = int((slide_width - img_w) / 2)
    img_top = 0

    prepared_image = get_prepared_image_path(image_path)

    if prepared_image:
        slide.shapes.add_picture(str(prepared_image), img_left, img_top, width=img_w, height=img_h)
    else:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            img_left,
            img_top,
            int(img_w),
            int(img_h),
        )
        box.text_frame.text = "Image not found"
        box.fill.background()
        box.line.fill.background()

    inset_left = int(img_w * TOP_DOWN_TEXT_LEFT_INSET_FRACTION)
    inset_right = int(img_w * TOP_DOWN_TEXT_RIGHT_INSET_FRACTION)

    text_left = img_left + inset_left
    text_top = int(img_top + img_h + text_gap)
    text_width = int(img_w - inset_left - inset_right)
    text_height = int(slide_height - text_top - bottom_margin)

    text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p1 = tf.paragraphs[0]
    p1.text = en_text
    p1.font.size = Pt(TOP_DOWN_EN_FONT_PT)
    p1.space_after = Pt(6)

    p2 = tf.add_paragraph()
    p2.text = zh_text
    p2.font.size = Pt(TOP_DOWN_ZH_FONT_PT)


def add_print4_card(
    slide,
    card_left: int,
    card_top: int,
    card_width: int,
    card_height: int,
    en_text: str,
    zh_text: str,
    image_path: Path | None,
) -> None:
    card_margin = Inches(PRINT4_CARD_MARGIN_IN)
    text_gap = Inches(PRINT4_TEXT_GAP_IN)

    usable_left = card_left + int(card_margin)
    usable_top = card_top + int(card_margin)
    usable_width = card_width - int(2 * card_margin)
    usable_height = card_height - int(2 * card_margin)

    img_w = usable_width
    img_h = int(usable_height * PRINT4_IMAGE_HEIGHT_FRACTION)

    prepared_image = get_prepared_image_path(image_path)

    if prepared_image:
        slide.shapes.add_picture(
            str(prepared_image),
            usable_left,
            usable_top,
            width=img_w,
            height=img_h,
        )
    else:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            usable_left,
            usable_top,
            img_w,
            img_h,
        )
        box.text_frame.text = "Image not found"
        box.fill.background()
        box.line.fill.background()

    inset_left = int(img_w * PRINT4_TEXT_LEFT_INSET_FRACTION)
    inset_right = int(img_w * PRINT4_TEXT_RIGHT_INSET_FRACTION)

    text_left = usable_left + inset_left
    text_top = usable_top + img_h + int(text_gap)
    text_width = img_w - inset_left - inset_right

    nominal_text_height = usable_height - img_h - int(text_gap)
    text_height = int(nominal_text_height * (1.0 - PRINT4_BOTTOM_RESERVE_FRACTION))

    fitted_en_pt, fitted_zh_pt = fit_font_sizes_for_two_paragraphs(
        en_text=en_text,
        zh_text=zh_text,
        box_width_emu=text_width,
        box_height_emu=text_height,
        start_en_pt=PRINT4_EN_FONT_PT,
        start_zh_pt=PRINT4_ZH_FONT_PT,
        paragraph_gap_pt=PRINT4_EN_AFTER_PT,
    )

    text_box = slide.shapes.add_textbox(
        text_left,
        text_top,
        text_width,
        text_height,
    )
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p1 = tf.paragraphs[0]
    p1.text = en_text
    p1.font.size = Pt(fitted_en_pt)
    p1.space_after = Pt(PRINT4_EN_AFTER_PT)

    p2 = tf.add_paragraph()
    p2.text = zh_text
    p2.font.size = Pt(fitted_zh_pt)


def build_presentation_left_right(rows: List[Dict[str, str]], media_dir: Path, output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(LEFT_RIGHT_SLIDE_WIDTH_IN)
    prs.slide_height = Inches(LEFT_RIGHT_SLIDE_HEIGHT_IN)

    for row in rows:
        img_name = row["ImgBaseName"].strip()
        image_path = media_dir / img_name if img_name else None
        add_slide_left_right(prs, row["En"], row["Zh"], image_path)

    prs.save(output_path)


def build_presentation_top_image_text_below(rows: List[Dict[str, str]], media_dir: Path, output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(TOP_DOWN_SLIDE_WIDTH_IN)
    prs.slide_height = Inches(TOP_DOWN_SLIDE_HEIGHT_IN)

    for row in rows:
        img_name = row["ImgBaseName"].strip()
        image_path = media_dir / img_name if img_name else None
        add_slide_top_image_text_below(prs, row["En"], row["Zh"], image_path)

    prs.save(output_path)


def build_presentation_print4(rows: List[Dict[str, str]], media_dir: Path, output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(PRINT4_PAGE_WIDTH_IN)
    prs.slide_height = Inches(PRINT4_PAGE_HEIGHT_IN)

    margin_left = Inches(PRINT4_PAGE_MARGIN_LEFT_IN)
    margin_right = Inches(PRINT4_PAGE_MARGIN_RIGHT_IN)
    margin_top = Inches(PRINT4_PAGE_MARGIN_TOP_IN)
    margin_bottom = Inches(PRINT4_PAGE_MARGIN_BOTTOM_IN)
    gap_x = Inches(PRINT4_CARD_GAP_X_IN)
    gap_y = Inches(PRINT4_CARD_GAP_Y_IN)

    usable_width = prs.slide_width - int(margin_left + margin_right)
    usable_height = prs.slide_height - int(margin_top + margin_bottom)

    card_width = int((usable_width - gap_x) / 2)
    card_height = int((usable_height - gap_y) / 2)

    positions = [
        (int(margin_left), int(margin_top)),
        (int(margin_left + card_width + gap_x), int(margin_top)),
        (int(margin_left), int(margin_top + card_height + gap_y)),
        (int(margin_left + card_width + gap_x), int(margin_top + card_height + gap_y)),
    ]

    for i in range(0, len(rows), 4):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        chunk = rows[i:i + 4]

        for j, row in enumerate(chunk):
            card_left, card_top = positions[j]
            img_name = row["ImgBaseName"].strip()
            image_path = media_dir / img_name if img_name else None

            add_print4_card(
                slide=slide,
                card_left=card_left,
                card_top=card_top,
                card_width=card_width,
                card_height=card_height,
                en_text=row["En"],
                zh_text=row["Zh"],
                image_path=image_path,
            )

    prs.save(output_path)


def convert_pptx_to_pdf(pptx_path: Path) -> Path | None:
    soffice = shutil.which("soffice")
    if not soffice:
        return None

    outdir = pptx_path.parent
    try:
        subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "pdf",
                str(pptx_path),
                "--outdir",
                str(outdir),
            ],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
    except subprocess.CalledProcessError:
        return None

    pdf_path = outdir / (pptx_path.stem + ".pdf")
    return pdf_path if pdf_path.exists() else None


def maybe_convert_to_pdf(paths: List[Path]) -> None:
    if not AUTO_CONVERT_TO_PDF:
        return

    for path in paths:
        pdf_path = convert_pptx_to_pdf(path)
        if pdf_path:
            print(f"PDF created: {pdf_path}")
        else:
            print(f"PDF conversion skipped/failed: {path}")


def prepare_rows(input_file: Path) -> List[Dict[str, str]]:
    rows = parse_anki_export(input_file)
    rows = dedupe_neighbor_overlaps(rows)
    return rows


def test_run() -> None:
    input_file = Path("02.txt")
    media_dir = DEFAULT_MEDIA_DIR

    rows = prepare_rows(input_file)
    export_sanitized_csv(rows, input_file.with_name("02_sanitized.csv"))

    print4_output = input_file.with_name(input_file.stem + "_print4_page.pptx")

    build_presentation_print4(rows, media_dir, print4_output)

    print(f"Rows parsed: {len(rows)}")
    print(f"CSV created: {input_file.with_name('02_sanitized.csv')}")
    print(f"PPTX created: {print4_output}")

    maybe_convert_to_pdf([print4_output])


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("input_file", help="Anki export text file")
    parser.add_argument("--media-dir", default=str(DEFAULT_MEDIA_DIR))
    parser.add_argument("--output-left-right", help="Output path for left/right pptx")
    parser.add_argument("--output-top-down", help="Output path for top/down pptx")
    parser.add_argument("--output-print4", help="Output path for 4-up letter-page pptx")
    parser.add_argument("--export-csv", action="store_true", help="Also export sanitized CSV")
    parser.add_argument("--no-pdf", action="store_true", help="Do not try to convert pptx files to pdf")
    args = parser.parse_args()

    global AUTO_CONVERT_TO_PDF
    if args.no_pdf:
        AUTO_CONVERT_TO_PDF = False

    input_path = Path(args.input_file).expanduser()
    media_dir = Path(args.media_dir).expanduser()

    left_right_output = (
        Path(args.output_left_right).expanduser()
        if args.output_left_right
        else input_path.with_suffix(".pptx")
    )
    top_down_output = (
        Path(args.output_top_down).expanduser()
        if args.output_top_down
        else input_path.with_name(input_path.stem + "_top_image.pptx")
    )
    print4_output = (
        Path(args.output_print4).expanduser()
        if args.output_print4
        else input_path.with_name(input_path.stem + "_print4_page.pptx")
    )

    rows = prepare_rows(input_path)

    if args.export_csv:
        export_sanitized_csv(rows, input_path.with_name(input_path.stem + "_sanitized.csv"))

    build_presentation_left_right(rows, media_dir, left_right_output)
    build_presentation_top_image_text_below(rows, media_dir, top_down_output)
    build_presentation_print4(rows, media_dir, print4_output)

    print(f"Rows parsed: {len(rows)}")
    print(f"PPTX created: {left_right_output}")
    print(f"PPTX created: {top_down_output}")
    print(f"PPTX created: {print4_output}")

    maybe_convert_to_pdf([left_right_output, top_down_output, print4_output])


if __name__ == "__main__":
    test_run()
    # main()